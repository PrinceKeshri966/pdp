import pypdf
import docx
from pptx import Presentation
import io

async def extract_text_from_file(file_content: bytes, filename: str) -> str:
    ext = filename.split('.')[-1].lower()
    text = ""
    stream = io.BytesIO(file_content)
    if ext == 'pdf':
        reader = pypdf.PdfReader(stream)
        text = "\n".join([page.extract_text() or "" for page in reader.pages])
    elif ext == 'docx':
        doc = docx.Document(stream)
        text = "\n".join([para.text for para in doc.paragraphs])
    elif ext == 'pptx':
        prs = Presentation(stream)
        text = "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text") and getattr(shape, "text")])
    return text
